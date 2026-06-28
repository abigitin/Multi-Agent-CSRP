from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs"

NAVY = PptColor(11, 22, 32)
BLUE = PptColor(23, 92, 211)
TEAL = PptColor(4, 120, 87)
GRAY = PptColor(82, 101, 118)
LIGHT = PptColor(245, 248, 251)


def add_textbox(slide, left, top, width, height, text, size=18, color=NAVY, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = PptPt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, bullets, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = PptPt(size)
        paragraph.font.color.rgb = GRAY
        paragraph.space_after = PptPt(7)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, PptInches(0.55), PptInches(0.35), PptInches(8.8), PptInches(0.45), title, 24, NAVY, True)
    if subtitle:
        add_textbox(slide, PptInches(0.58), PptInches(0.83), PptInches(8.2), PptInches(0.32), subtitle, 11, GRAY, False)


def add_card(slide, left, top, width, height, title, body, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = PptColor(217, 226, 234)
    shape.adjustments[0] = 0.08
    add_textbox(slide, left + PptInches(0.16), top + PptInches(0.12), width - PptInches(0.32), PptInches(0.25), title, 13, accent, True)
    add_textbox(slide, left + PptInches(0.16), top + PptInches(0.43), width - PptInches(0.32), height - PptInches(0.52), body, 11, GRAY, False)


def add_arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = PptColor(130, 146, 161)
    line.line.width = PptPt(1.5)
    return line


def generate_pptx() -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PptColor(238, 242, 245)
    add_textbox(slide, PptInches(0.7), PptInches(1.25), PptInches(8.8), PptInches(0.38), "AgenticAI SupportOps", 16, BLUE, True)
    add_textbox(slide, PptInches(0.68), PptInches(1.75), PptInches(9.3), PptInches(1.25), "Customer Support Resolution Platform", 42, NAVY, True)
    add_textbox(
        slide,
        PptInches(0.72),
        PptInches(3.1),
        PptInches(8.8),
        PptInches(0.7),
        "A multi-agent workflow for ticket intake, Pinecone knowledge retrieval, AI drafting, guardrails, approval, and Gmail SMTP delivery.",
        18,
        GRAY,
        False,
    )
    add_card(slide, PptInches(9.55), PptInches(1.38), PptInches(2.85), PptInches(1.1), "Live RAG", "Pinecone customer-support\nnamespace __default__", TEAL)
    add_card(slide, PptInches(9.55), PptInches(2.75), PptInches(2.85), PptInches(1.1), "SMTP", "Gmail notification path\nverified sent", BLUE)
    add_card(slide, PptInches(9.55), PptInches(4.12), PptInches(2.85), PptInches(1.1), "Frontend", "React operations console\nprofessional dashboard", NAVY)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "What Was Built", "Business-friendly summary of the platform")
    add_bullets(
        slide,
        PptInches(0.7),
        PptInches(1.45),
        PptInches(5.7),
        PptInches(4.9),
        [
            "A React dashboard for support operators.",
            "A FastAPI backend that manages tickets, runs, drafts, citations, approvals, and notifications.",
            "A Pinecone integrated-embedding knowledge store for live semantic retrieval.",
            "A multi-agent workflow that records every step for review and audit.",
            "A Gmail SMTP notification path that can send through the app.",
        ],
        18,
    )
    add_card(slide, PptInches(7.0), PptInches(1.5), PptInches(2.45), PptInches(1.2), "Operator", "Reviews tickets and approves responses", BLUE)
    add_card(slide, PptInches(9.75), PptInches(1.5), PptInches(2.45), PptInches(1.2), "System", "Runs agents and stores the full audit trail", TEAL)
    add_card(slide, PptInches(7.0), PptInches(3.2), PptInches(5.2), PptInches(1.35), "Outcome", "Faster ticket response with citations, guardrails, and delivery tracking.", NAVY)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "High-Level Design", "How the main building blocks talk to each other")
    boxes = [
        ("Operator", 0.6, 1.25, "Support user works in the dashboard"),
        ("React Frontend", 2.7, 1.25, "Ticket dashboard and review UI"),
        ("FastAPI Backend", 5.0, 1.25, "REST API and orchestration"),
        ("Agent Workflow", 7.35, 1.25, "Intake, retrieval, draft, guardrails"),
        ("MCP Tools", 9.75, 1.25, "Knowledge, memory, mail, ticket tools"),
        ("SQLite", 5.0, 3.55, "Runs, traces, drafts, approvals"),
        ("Pinecone", 8.0, 3.55, "Integrated embedding knowledge search"),
        ("Gmail SMTP", 10.85, 3.55, "Customer notification delivery"),
    ]
    coords = {}
    for title, x, y, body in boxes:
        left = PptInches(x)
        top = PptInches(y)
        width = PptInches(1.85)
        height = PptInches(1.0)
        add_card(slide, left, top, width, height, title, body, BLUE if title in {"React Frontend", "FastAPI Backend"} else TEAL)
        coords[title] = (left, top, width, height)
    for first, second in [
        ("Operator", "React Frontend"),
        ("React Frontend", "FastAPI Backend"),
        ("FastAPI Backend", "Agent Workflow"),
        ("Agent Workflow", "MCP Tools"),
    ]:
        x1, y1, w1, h1 = coords[first]
        x2, y2, _, h2 = coords[second]
        add_arrow(slide, x1 + w1, y1 + h1 / 2, x2, y2 + h2 / 2)
    add_arrow(slide, coords["FastAPI Backend"][0] + PptInches(0.92), coords["FastAPI Backend"][1] + PptInches(1.0), coords["SQLite"][0] + PptInches(0.92), coords["SQLite"][1])
    add_arrow(slide, coords["MCP Tools"][0] + PptInches(0.3), coords["MCP Tools"][1] + PptInches(1.0), coords["Pinecone"][0] + PptInches(0.9), coords["Pinecone"][1])
    add_arrow(slide, coords["MCP Tools"][0] + PptInches(1.55), coords["MCP Tools"][1] + PptInches(1.0), coords["Gmail SMTP"][0] + PptInches(0.8), coords["Gmail SMTP"][1])

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Agent Workflow", "Each step has a clear responsibility and leaves an audit trail")
    steps = [
        ("1. Intake", "Normalize ticket and issue category"),
        ("2. Retrieval", "Search Pinecone for relevant knowledge"),
        ("3. Resolution", "Draft a customer response"),
        ("4. Communication", "Send or queue notification"),
        ("5. Evaluation", "Score confidence"),
        ("6. Guardrails", "Pass or route to review"),
    ]
    for index, (title, body) in enumerate(steps):
        row = index // 3
        col = index % 3
        add_card(
            slide,
            PptInches(0.75 + col * 4.1),
            PptInches(1.55 + row * 2.0),
            PptInches(3.55),
            PptInches(1.25),
            title,
            body,
            BLUE if index < 3 else TEAL,
        )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Key Files", "Where the important behavior lives")
    add_bullets(
        slide,
        PptInches(0.7),
        PptInches(1.35),
        PptInches(5.9),
        PptInches(5.4),
        [
            "backend/main.py starts the API and database.",
            "backend/api/routes.py exposes frontend endpoints.",
            "backend/agents/workflow.py runs the agent process.",
            "backend/rag/pinecone_store.py handles Pinecone upsert and search.",
            "backend/integrations/mail.py sends SMTP notifications.",
        ],
        15,
    )
    add_bullets(
        slide,
        PptInches(7.0),
        PptInches(1.35),
        PptInches(5.6),
        PptInches(5.4),
        [
            "frontend/src/components/TicketDashboard.tsx is the dashboard screen.",
            "frontend/src/api.ts calls the backend.",
            "frontend/src/styles.css contains the new UI design.",
            "README.md explains the full app.",
            "AGENTS.md explains agent responsibilities.",
        ],
        15,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Verification Results", "What was confirmed after implementation")
    add_bullets(
        slide,
        PptInches(0.8),
        PptInches(1.35),
        PptInches(11.7),
        PptInches(4.8),
        [
            "Frontend production build passes.",
            "Backend tests pass.",
            "Pinecone status reports live.",
            "Pinecone namespace __default__ contains synced vectors.",
            "MCP knowledge search returns pinecone mode results.",
            "Ticket workflow persists retrieval_mode=pinecone with citations.",
            "Gmail SMTP app-path notification returned sent.",
        ],
        19,
    )

    prs.save(OUT / "AgenticAI_HLD_and_Project_Overview.pptx")


def add_doc_heading(doc: Document, text: str, level: int = 1):
    heading = doc.add_heading(text, level=level)
    for run in heading.runs:
        run.font.color.rgb = RGBColor(11, 22, 32)
    return heading


def add_doc_bullets(doc: Document, items: list[str]) -> None:
    for item in items:
        doc.add_paragraph(item, style="List Bullet")


def generate_docx() -> None:
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("AgenticAI Customer Support Resolution Platform")
    run.bold = True
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(11, 22, 32)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.add_run("Plain-English technical explanation, HLD, file guide, and operating plan").font.size = Pt(11)

    add_doc_heading(doc, "1. What This Application Does")
    doc.add_paragraph(
        "This application helps a support team resolve customer tickets by combining a React dashboard, "
        "a FastAPI backend, Pinecone knowledge search, a multi-agent workflow, and Gmail SMTP notifications."
    )
    add_doc_bullets(
        doc,
        [
            "Tickets come from ServiceNow or local demo data.",
            "Knowledge comes from Atlassian or local markdown files.",
            "Pinecone stores searchable knowledge using integrated embeddings.",
            "Agents retrieve evidence, draft a response, evaluate confidence, apply guardrails, and record every step.",
            "Operators review the response, approve it, and track notification delivery.",
        ],
    )

    add_doc_heading(doc, "2. High-Level Design")
    doc.add_paragraph(
        "Operator -> React Frontend -> FastAPI Backend -> Agent Workflow -> MCP Tools -> Pinecone and Gmail SMTP. "
        "SQLite stores tickets, runs, traces, drafts, citations, approvals, memory, and notifications."
    )
    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    hdr[0].text = "Layer"
    hdr[1].text = "Technology"
    hdr[2].text = "Responsibility"
    rows = [
        ("Frontend", "React + Vite", "Operator dashboard and review workspace"),
        ("Backend API", "FastAPI", "REST endpoints and orchestration"),
        ("Agent Workflow", "Python workflow / LangGraph compatible", "Ticket resolution steps"),
        ("Tool Layer", "MCP server/client", "Knowledge, ticket, memory, and notification tools"),
        ("Knowledge Search", "Pinecone", "Integrated-embedding semantic retrieval"),
        ("Notification", "Gmail SMTP", "Customer email delivery"),
        ("Database", "SQLite", "Audit trail and application state"),
    ]
    for row in rows:
        cells = table.add_row().cells
        for index, value in enumerate(row):
            cells[index].text = value

    add_doc_heading(doc, "3. How A Ticket Run Works")
    add_doc_bullets(
        doc,
        [
            "The operator clicks Run Agents.",
            "The frontend calls the backend run endpoint.",
            "The backend loads the ticket and starts the agent workflow.",
            "The retrieval agent searches Pinecone using the ticket text.",
            "The resolution agent creates a draft using the retrieved evidence.",
            "The communication agent sends or queues the notification.",
            "The evaluation and guardrail agents decide whether review is needed.",
            "All outputs are stored and shown on the dashboard.",
        ],
    )

    add_doc_heading(doc, "4. File-By-File Explanation")
    files = [
        ("backend/main.py", "Starts the FastAPI app, enables CORS, validates settings, and initializes the database."),
        ("backend/api/routes.py", "Defines the API endpoints used by the frontend."),
        ("backend/core/config.py", "Loads .env settings for Pinecone, SMTP, LLM, ServiceNow, Atlassian, MCP, and database."),
        ("backend/core/database.py", "Creates the database connection and initializes/migrates tables."),
        ("backend/core/models.py", "Defines database tables."),
        ("backend/core/status.py", "Builds the provider status cards shown in the UI."),
        ("backend/agents/workflow.py", "Runs the multi-agent ticket workflow."),
        ("backend/agents/llm.py", "Calls the LLM or uses fallback response generation."),
        ("backend/mcp/server.py", "Implements tools used by the agents."),
        ("backend/rag/knowledge_sync.py", "Syncs knowledge into SQLite and Pinecone."),
        ("backend/rag/pinecone_store.py", "Central Pinecone integrated-embedding upsert and search logic."),
        ("backend/integrations/mail.py", "Sends Gmail SMTP notifications."),
        ("frontend/src/components/TicketDashboard.tsx", "Main dashboard screen."),
        ("frontend/src/api.ts", "Frontend API client."),
        ("frontend/src/styles.css", "Professional UI styling and responsive layout."),
    ]
    table = doc.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    table.rows[0].cells[0].text = "File"
    table.rows[0].cells[1].text = "Purpose"
    for path, purpose in files:
        cells = table.add_row().cells
        cells[0].text = path
        cells[1].text = purpose

    add_doc_heading(doc, "5. Agent Responsibilities")
    agent_rows = [
        ("Intake", "Normalizes ticket context."),
        ("Retrieval", "Finds citations from Pinecone."),
        ("Resolution", "Creates the customer draft."),
        ("Communication", "Sends or queues email."),
        ("Evaluation", "Scores confidence."),
        ("Guardrails", "Checks quality and safety."),
        ("Approval", "Completes or routes to human review."),
    ]
    for name, purpose in agent_rows:
        paragraph = doc.add_paragraph()
        paragraph.add_run(f"{name}: ").bold = True
        paragraph.add_run(purpose)

    add_doc_heading(doc, "6. Verification Results")
    add_doc_bullets(
        doc,
        [
            "Frontend build passes.",
            "Backend tests pass.",
            "Pinecone is live and uses namespace __default__.",
            "Knowledge sync upserts records to Pinecone.",
            "Ticket workflow returns retrieval_mode=pinecone.",
            "SMTP notification sends successfully through the app path.",
        ],
    )

    add_doc_heading(doc, "7. Important Non-Coder Notes")
    add_doc_bullets(
        doc,
        [
            "Use the dashboard buttons; you do not need to run code for daily use.",
            "Keep .env private because it contains keys and passwords.",
            "If email stops working, create a new Gmail app password and update SMTP_PASSWORD.",
            "If Pinecone search stops working, check the Pinecone key, host, index, and namespace.",
            "Mock data keeps the app usable without ServiceNow or Atlassian credentials.",
        ],
    )

    doc.save(OUT / "AgenticAI_Project_Explanation.docx")


if __name__ == "__main__":
    OUT.mkdir(exist_ok=True)
    generate_pptx()
    generate_docx()
    print("Generated docs/AgenticAI_HLD_and_Project_Overview.pptx")
    print("Generated docs/AgenticAI_Project_Explanation.docx")
